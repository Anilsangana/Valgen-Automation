"""
Generate Professional PowerPoint Presentation for DLWA
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Color scheme
PRIMARY_BLUE = RGBColor(0, 82, 165)  # #0052A5
DARK_BLUE = RGBColor(0, 61, 122)     # #003D7A
ACCENT_BLUE = RGBColor(0, 163, 224)  # #00A3E0
TEXT_DARK = RGBColor(45, 55, 72)     # #2D3748
TEXT_GRAY = RGBColor(113, 128, 150)  # #718096
SUCCESS_GREEN = RGBColor(72, 187, 120)
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add blue gradient background (simulated with rectangle)
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    
    bg = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY_BLUE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Digital Lifecycle Workflow Accelerator"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Automated User Lifecycle Management System"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = WHITE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Badges
    badges_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(0.5))
    badges_frame = badges_box.text_frame
    badges_frame.text = "GCP Compliant  |  21 CFR Part 11  |  100% Audit Coverage"
    badges_para = badges_frame.paragraphs[0]
    badges_para.font.size = Pt(16)
    badges_para.font.color.rgb = ACCENT_BLUE
    badges_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    """Add a content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_para = title_shape.text_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE
    
    # Content
    content = slide.shapes.placeholders[1]
    text_frame = content.text_frame
    text_frame.clear()
    
    for bullet in bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_DARK

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)

def add_metrics_slide(prs, title, metrics):
    """Add a slide with large metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE
    
    # Metrics boxes
    x_start = 1
    y_pos = 2
    box_width = 2.5
    
    for i, (metric_value, metric_label) in enumerate(metrics):
        x_pos = x_start + (i * (box_width + 0.3))
        
        # Metric value
        value_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos), Inches(box_width), Inches(1))
        value_frame = value_box.text_frame
        value_frame.text = metric_value
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(48)
        value_para.font.bold = True
        value_para.font.color.rgb = PRIMARY_BLUE
        value_para.alignment = PP_ALIGN.CENTER
        
        # Metric label
        label_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_pos + 1.2), Inches(box_width), Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = metric_label
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(16)
        label_para.font.color.rgb = TEXT_GRAY
        label_para.alignment = PP_ALIGN.CENTER

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs)
    
    # Slide 2: Executive Summary
    add_content_slide(prs, "Executive Summary", [
        "Professional automation platform for enterprise validation workflows",
        "Automated User Management - Reduce manual effort by 90%",
        "GCP Compliance - Built-in audit trails",
        "21 CFR Part 11 Validation - Pharmaceutical-grade compliance",
        "100% Audit Coverage - Every action logged and documented"
    ])
    
    # Slide 3: Problem Statement
    add_content_slide(prs, "Traditional Challenges", [
        "Manual user creation takes 15-20 minutes per user",
        "High risk of human error in data entry",
        "GCP audits require detailed activity logs",
        "PDF documentation must be created manually",
        "No standardized audit format",
        "Admin spends hours on repetitive tasks"
    ])
    
    # Slide 4: Our Solution
    add_content_slide(prs, "Our Solution", [
        "Single dashboard to manage all administrative tasks",
        "Automated role, department, and user provisioning",
        "Built-in approval workflows with real-time logging",
        "Automatic PDF audit trail generation",
        "GCP-compliant documentation",
        "Enterprise-grade quality and security"
    ])
    
    # Slide 5: Core Features
    add_content_slide(prs, "Core Features Overview", [
        "Role Management - Create custom roles with permissions",
        "Department Management - Automated setup workflow",
        "User Provisioning - Full lifecycle automation",
        "User Deactivation - Secure deactivation process",
        "Complete Setup Flow - End-to-end automation",
        "PDF Audit Trails - Regulatory compliance built-in"
    ])
    
    # Slide 6: Role Creation
    add_content_slide(prs, "Role Management Workflow", [
        "Automated permission assignment",
        "Duplicate role detection with flexible strategies",
        "Real-time status logging and validation",
        "Error recovery mechanisms",
        "PDF audit trail with timestamps",
        "Time Saved: 5-7 minutes per role"
    ])
    
    # Slide 7: Department Creation
    add_content_slide(prs, "Department Management", [
        "Hierarchical organization support",
        "Automated description population",
        "Duplicate detection and handling",
        "Batch processing capability",
        "Complete audit documentation",
        "Time Saved: 3-5 minutes per department"
    ])
    
    # Slide 8: User Creation
    add_content_slide(prs, "User Provisioning - 4-Step Process", [
        "Step 1: Admin Form Submission - Automated form filling",
        "Step 2: User Signup - Email verification simulation",
        "Step 3: Admin Activation - Automated approval workflow",
        "Step 4: Login Verification - Credentials validation",
        "Every step logged with precise timestamp",
        "Time Saved: 15-20 minutes per user"
    ])
    
    # Slide 9: User Deactivation
    add_content_slide(prs, "User Deactivation Process", [
        "Batch deactivation support",
        "Admin-authenticated process only",
        "Continue-on-error option for bulk operations",
        "Detailed status messages and logging",
        "Deactivation audit trail with timestamps",
        "Secure compliance documentation"
    ])
    
    # Slide 10: Complete Setup Flow
    add_content_slide(prs, "Complete Setup Flow", [
        "Single Click → Complete Setup",
        "Role Creation → Department Setup → User Provisioning",
        "Comprehensive audit report for entire workflow",
        "30-40 minutes saved per complete setup",
        "Zero manual intervention required",
        "Consistent process every time"
    ])
    
    # Slide 11: Professional UI
    add_content_slide(prs, "Professional User Interface", [
        "Modern design with DLWA branding",
        "GCP Compliant and 21 CFR Part 11 badges",
        "Clear value proposition with key statistics",
        "Real-time activity logging with color coding",
        "Easy-to-use form inputs and controls",
        "Live system status indicator"
    ])
    
    # Slide 12: PDF Audit Trail
    add_content_slide(prs, "GCP Compliance & Audit Trail", [
        "Automatically generated for every operation",
        "Professional branded header in enterprise blue",
        "Operation details: Type, User, Environment, Timestamp",
        "Execution results with individual timestamps",
        "Page footers with Report ID and page numbers",
        "GCP compliance statement on every page"
    ])
    
    # Slide 13: PDF Report Structure
    add_content_slide(prs, "PDF Audit Report Features", [
        "Blue branded header - Professional appearance",
        "Operation Details - Complete context",
        "Execution Results - Item-by-item documentation",
        "Individual timestamps (IST) - Full traceability",
        "Status tracking - created/skipped/failed",
        "Unique Report ID - Easy retrieval"
    ])
    
    # Slide 14: Compliance Benefits
    add_content_slide(prs, "Meeting Regulatory Requirements", [
        "GCP: Audit trail, traceability, accountability",
        "21 CFR Part 11: Electronic records & signatures",
        "Data Integrity: Immutable PDF records",
        "Version Control: Unique report IDs",
        "Standardized format: Consistent structure",
        "Inspection ready: Regulator-friendly format"
    ])
    
    # Slide 15: Performance Metrics
    add_metrics_slide(prs, "Performance & Time Savings", [
        ("90%", "Time Saved"),
        ("100+ hrs", "Annually"),
        ("$10K", "Cost Reduction")
    ])
    
    # Slide 16: Reliability
    add_content_slide(prs, "Enterprise-Grade Reliability", [
        "Optimized performance - No unnecessary session checks",
        "Automatic retry on transient failures",
        "Duplicate detection with flexible strategies",
        "Network error recovery mechanisms",
        "95%+ success rate for user creation",
        "100% audit trail generation"
    ])
    
    # Slide 17: Use Case - Project Onboarding
    add_two_column_slide(prs, 
        "Use Case: New Project Onboarding",
        [
            "Traditional Approach:",
            "• 10 users × 20 min = 200 min",
            "• Manual PDF = 30 min",
            "• Total: 230 minutes",
            "",
            "Error rate: ~15%",
            "Compliance: Manual"
        ],
        [
            "With DLWA:",
            "• 10 users × 3 min = 30 min",
            "• Automatic PDF = Instant",
            "• Total: 30 minutes",
            "",
            "Error rate: 0%",
            "Compliance: Automatic"
        ]
    )
    
    # Slide 18: ROI Calculation
    add_content_slide(prs, "Return on Investment", [
        "Annual Time Savings: 100+ hours",
        "Cost Savings: $5,000 - $10,000 per year",
        "Error Reduction: 90%+ fewer manual errors",
        "Compliance Value: Invaluable during audits",
        "Payback Period: Weeks, not months",
        "Ongoing savings compound over time"
    ])
    
    # Slide 19: Key Differentiators
    add_content_slide(prs, "Why Choose DLWA?", [
        "Complete Automation - Full lifecycle, not just creation",
        "Regulatory Compliance - GCP & 21 CFR Part 11 built-in",
        "Professional Quality - Enterprise-grade everything",
        "Proven Reliability - 95%+ success rate",
        "Time & Cost Savings - 90% reduction in effort",
        "Modern Technology - Future-proof architecture"
    ])
    
    # Slide 20: Security Features
    add_content_slide(prs, "Enterprise Security Standards", [
        "Admin authentication required for all operations",
        "Role-based access control (RBAC)",
        "Secure session management",
        "Complete audit logging and accountability",
        "Immutable PDF records - Data integrity",
        "Timestamp precision with IST timezone"
    ])
    
    # Slide 21: Technical Architecture
    add_content_slide(prs, "Modern, Scalable Design", [
        "Frontend: Modern HTML5 + CSS3 + JavaScript",
        "Backend: Node.js with TypeScript",
        "Automation: Playwright browser automation",
        "PDF: Professional PDFKit library",
        "Modular architecture - Clean code structure",
        "Enterprise error handling throughout"
    ])
    
    # Slide 22: Implementation
    add_content_slide(prs, "Getting Started", [
        "Cloud or On-Premise deployment options",
        "2-hour onboarding and training session",
        "Complete documentation provided",
        "Email and phone support available",
        "Monthly feature updates",
        "Dedicated account manager (Professional+)"
    ])
    
    # Slide 23: Demo Overview
    add_content_slide(prs, "Live Demonstration", [
        "UI Walkthrough - Professional interface",
        "Role Creation - Quick automation demo",
        "Complete Setup Flow - End-to-end process",
        "PDF Review - Audit trail examination",
        "Real-time logging - Activity monitoring",
        "Q&A - Your specific questions"
    ])
    
    # Slide 24: Pricing
    add_content_slide(prs, "Flexible Pricing Options", [
        "Starter: $499/month - Up to 50 users/month",
        "Professional: $1,299/month - Up to 200 users/month",
        "Enterprise: Custom pricing - Unlimited users",
        "All plans include PDF audit trails",
        "Free 14-day trial available",
        "30-day pilot program option"
    ])
    
    # Slide 25: Next Steps
    add_content_slide(prs, "Let's Get Started!", [
        "Schedule a Live Demo - See it in action",
        "Free Trial - 14 days, no credit card required",
        "Pilot Program - 30-day evaluation",
        "Contact: sales@dlwa-automation.com",
        "Phone: +1 (555) 123-4567",
        "Website: www.dlwa-automation.com"
    ])
    
    # Slide 26: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Questions & Discussion"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY_BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "We're here to help!"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = TEXT_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('DLWA_Client_Presentation.pptx')
    print("✅ PowerPoint presentation created successfully!")
    print("📄 File: DLWA_Client_Presentation.pptx")
    print(f"📊 Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
